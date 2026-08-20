#!/usr/bin/env python3
"""Export a PPTD project to a flattened PPTX using Kimi-rendered page images."""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Optional, Sequence

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from export_images import export_images
from export_pptx import ExportError, find_manifest, patch_transitions, verify_output


def export_pptx(
    source: Path,
    output: Path,
    transition: str,
    force: bool,
) -> dict[str, object]:
    manifest = find_manifest(source)
    output = output.expanduser().resolve()
    if output.exists() and not force:
        raise ExportError(f"output already exists (pass --force): {output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    manifest_data = yaml.safe_load(manifest.read_text(encoding="utf-8"))
    title = str(manifest_data.get("title") or manifest.stem)

    with tempfile.TemporaryDirectory(prefix="open-kimi-ppt-local-") as temp_name:
        image_output = Path(temp_name) / "images"
        image_summary = export_images(manifest, image_output)
        image_paths = [image_output / item["image"] for item in image_summary["images"]]

        with Image.open(image_paths[0]) as first_image:
            width, height = first_image.size

        presentation = Presentation()
        presentation.slide_width = Inches(13.333333)
        presentation.slide_height = round(presentation.slide_width * height / width)
        blank_layout = presentation.slide_layouts[6]
        for image_path in image_paths:
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

        presentation.core_properties.title = title
        presentation.save(output)

    patched_slides = patch_transitions(output, transition)
    summary = verify_output(output, transition, expect_fonts=False)
    summary.update(
        {
            "flattened": True,
            "output": str(output),
            "transitionPatchedSlides": patched_slides,
        }
    )
    return summary


def parse_args(argv: Optional[Sequence[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Export PPTD to a flattened PPTX using Kimi-rendered page images.")
    parser.add_argument("input", type=Path, help=".pptd manifest or project directory")
    parser.add_argument("--output", "-o", type=Path, help="output .pptx path")
    parser.add_argument(
        "--transition",
        choices=("fade", "none"),
        default="fade",
        help="slide transition (default: fade)",
    )
    parser.add_argument("--force", action="store_true", help="replace an existing output file")
    return parser.parse_args(argv)


def main(argv: Optional[Sequence[str]] = None) -> int:
    args = parse_args(argv)
    try:
        manifest = find_manifest(args.input)
        output = args.output or manifest.with_suffix(".pptx")
        summary = export_pptx(args.input, output, args.transition, args.force)
    except (ExportError, OSError, subprocess.SubprocessError, ValueError) as exc:
        print(f"open-kimi-ppt local export failed: {exc}", file=sys.stderr)
        return 1
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
